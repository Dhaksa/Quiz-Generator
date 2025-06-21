import os
import faiss
import PyPDF2
import docx
from pptx import Presentation
import requests
from typing import List
from sentence_transformers import SentenceTransformer
from models.progress_tracker import log_score, update_xp_streak
from datetime import date

GROQ_API_KEY = "your_groq_cloud_api_key"  # Replace with your actual key
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL = "llama3-8b-8192"
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")


def read_document(file_path: str) -> str:
    text = ""
    file_path_lower = file_path.lower()
    if file_path_lower.endswith(".pdf"):
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    elif file_path_lower.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    elif file_path_lower.endswith(".docx"):
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_path_lower.endswith(".pptx"):
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_path_lower.endswith(".md"):
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    else:
        raise ValueError("Unsupported file type!")
    return text


def chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size - overlap)]


def build_faiss_index(chunks: List[str]):
    embeddings = embedding_model.encode(chunks)
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    return index, embeddings, chunks


def query_faiss(query: str, index, chunks, embeddings, top_k=3):
    query_embedding = embedding_model.encode([query])
    _, indices = index.search(query_embedding, top_k)
    return " ".join([chunks[i] for i in indices[0]])


def generate_quiz_list_with_groq(context: str, num_questions: int = 5) -> List[str]:
    prompt = f"""You are a quiz generator.

From the given context below, generate {num_questions} multiple-choice questions (MCQs).
Each question must have:
- A clear question
- Four options (A, B, C, D)
- The correct option (e.g., "Answer: B")

Return them in a numbered list format like:

1. Question...
   A. ...
   B. ...
   C. ...
   D. ...
   Answer: X

Context:
\"\"\"
{context}
\"\"\"
"""
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": GROQ_MODEL,
        "messages": [
            {"role": "system", "content": "You are a helpful quiz generator."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.9
    }

    response = requests.post(GROQ_API_URL, headers=headers, json=payload)

    if response.status_code == 200:
        content = response.json()["choices"][0]["message"]["content"]
        return content.split("\n\n")
    else:
        return [f"Error: {response.status_code} - {response.text}"]


def generate_cloze_test(context: str, num_questions=3) -> List[str]:
    prompt = f"""Generate {num_questions} fill-in-the-blank questions (cloze test) from the context below.
Format:
1. ____ is the capital of France. (Answer: Paris)

Context:
{context}
"""
    return call_groq_api(prompt).split("\n\n")


def generate_flashcards(context: str, num_cards=5) -> List[str]:
    prompt = f"""Extract {num_cards} important concept-definition flashcards from the following content.
Format:
Term: X
Definition: Y

Content:
{context}
"""
    return call_groq_api(prompt).split("\n\n")


def call_groq_api(prompt):
    headers = {"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"}
    payload = {
        "model": GROQ_MODEL,
        "messages": [
            {"role": "system", "content": "You are a helpful AI."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.9
    }
    response = requests.post(GROQ_API_URL, headers=headers, json=payload)
    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        return f"Error: {response.status_code}"


def parse_question_block(block: str):
    lines = block.strip().split("\n")
    question_text = []
    options = {}
    answer = None
    for line in lines:
        if line.strip().startswith("Answer:"):
            answer = line.strip().split("Answer:")[-1].strip().upper()
        elif line.strip()[0:2] in ("A.", "B.", "C.", "D."):
            key = line.strip()[0]
            value = line.strip()[2:].strip()
            options[key] = value
        else:
            question_text.append(line)
    return " ".join(question_text), options, answer


def run_quiz(questions: List[str]):
    score = 0
    total = 0
    for block in questions:
        if "Answer:" not in block:
            continue
        total += 1
        question, options, correct = parse_question_block(block)
        print("\nQuestion:", question)
        for key in sorted(options.keys()):
            print(f"   {key}. {options[key]}")
        user_ans = input("Your answer (A/B/C/D): ").strip().upper()
        if user_ans == correct:
            print("Correct!")
            score += 1
        else:
            print(f"Incorrect! Correct answer: {correct}")
    print(f"\nFinal Score: {score}/{total}")
    log_score(score, username)
    update_xp_streak(date.today(), score, username)
